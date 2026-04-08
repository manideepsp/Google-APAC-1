from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Pt


ROOT = Path(__file__).resolve().parent
PPT_PATH = ROOT / "HackathonSubmission.pptx"


def set_text(shape, text: str, font_size: int = 14) -> None:
    tf = shape.text_frame
    tf.clear()

    lines = [line.rstrip() for line in text.strip().splitlines() if line.strip()]
    if not lines:
        return

    tf.paragraphs[0].text = lines[0]
    for line in lines[1:]:
        tf.add_paragraph().text = line

    for paragraph in tf.paragraphs:
        paragraph.space_after = Pt(2)
        for run in paragraph.runs:
            run.font.size = Pt(font_size)

    tf.word_wrap = True


def clear_text(shape) -> None:
    shape.text_frame.clear()


def add_picture_fit(slide, image_path: Path, left: int, top: int, width: int, height: int) -> None:
    with Image.open(image_path) as img:
        img_w, img_h = img.size

    box_w = float(width)
    box_h = float(height)
    scale = min(box_w / img_w, box_h / img_h)

    pic_w = int(img_w * scale)
    pic_h = int(img_h * scale)
    pic_left = int(left + (box_w - pic_w) / 2)
    pic_top = int(top + (box_h - pic_h) / 2)

    slide.shapes.add_picture(str(image_path), pic_left, pic_top, width=pic_w, height=pic_h)


def add_textbox(slide, left: int, top: int, width: int, height: int, text: str, font_size: int = 11) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    set_text(box, text, font_size=font_size)


def main() -> None:
    prs = Presentation(str(PPT_PATH))

    # Slide 1
    s1 = prs.slides[0]
    set_text(
        s1.shapes[1],
        "Participant Details\n"
        "Participant Name: Manideep SP\n"
        "Problem Statement: Build an AI strategy console that converts a creator's goal into a practical, trackable, and improvable execution plan using ADK agents + MCP tools.\n"
        "Hackathon Focus: Goal-to-execution workflow with live task lifecycle, analytics-driven reflection, and AlloyDB AI readiness.",
        font_size=18,
    )

    # Slide 2
    s2 = prs.slides[1]
    set_text(
        s2.shapes[2],
        "AI Strategy Console is a multi-agent YouTube strategy copilot designed for real creator workflows, not one-shot idea generation.\n"
        "- User provides goal, audience, budget, timeline, and channel context.\n"
        "- ADK agents perform research, planning, execution, and reflection in a closed loop.\n"
        "- MCP tools connect external signals and operations: YouTube trends/analytics, web search, SQLite task state, and Sheets sync.\n"
        "- Output is a day-wise action plan with traceable edits, run history, and re-analysis for continuous improvement.\n"
        "This delivers immediate productivity for creators today while keeping a clean migration path to AlloyDB AI for scale.",
        font_size=14,
    )

    # Slide 3
    s3 = prs.slides[2]
    set_text(
        s3.shapes[1],
        "Core approach/workflow (ADK + MCP):\n"
        "1. Capture structured goal parameters from user input.\n"
        "2. Research Agent pulls YouTube trends, analytics, web insights, and KPIs.\n"
        "3. Planning Agent produces prioritized day-wise task plans.\n"
        "4. Execution Agent saves tasks to SQLite and syncs active tasks to Sheets.\n"
        "5. User executes/edits tasks with traceable status and live/archive lifecycle.\n"
        "6. Reflection Agent re-analyzes performance and regenerates improved tasks.\n"
        "Loop outcome: Goal -> Research -> Plan -> Execute -> Reflect -> Replan.",
        font_size=10,
    )

    # Slide 4
    s4 = prs.slides[3]
    set_text(
        s4.shapes[2],
        "Differentiation from existing ideas:\n"
        "- Not just an idea generator: it closes the loop from strategy to execution and reflection.\n"
        "- Combines ADK reasoning with MCP tool control inside one operational workflow.\n"
        "- Captures task change traces (field-level updates and move history) for accountability.\n"
        "- Supports per-user, per-channel run history instead of one-off prompts.\n"
        "USP:\n"
        "- End-to-end Goal -> Research -> Plan -> Execute -> Reflect system with production-like boundaries.\n"
        "- Practical hybrid design: fast local SQLite now, phased AlloyDB AI path for analytics and RAG later.",
        font_size=14,
    )

    # Slide 5
    s5 = prs.slides[4]
    set_text(
        s5.shapes[2],
        "- Secure auth flow: register/login/logout/reset with PBKDF2 password hashing.\n"
        "- Conversational goal assistant that iteratively builds strategy parameters.\n"
        "- ADK multi-agent orchestration for research, planning, execution, and reflection.\n"
        "- KPI builder combining trending and channel-performance insights.\n"
        "- Task workspace with priority/status/live lifecycle controls.\n"
        "- Run history, archive, and task-modification trace feed.\n"
        "- YouTube and Sheets gRPC microservices for external integrations.\n"
        "- MCP server exposing 6 reusable operational tools.\n"
        "- SQLite WAL + retry-safe updates for reliable local persistence.\n"
        "- Sync of active tasks to Google Sheets for collaboration.",
        font_size=13,
    )

    # Slide 6: process flow diagram
    s6 = prs.slides[5]
    set_text(s6.shapes[2], "")
    box = s6.shapes[2]
    add_picture_fit(
        s6,
        ROOT / "diagrams" / "process-flow.png",
        box.left,
        box.top,
        box.width,
        box.height,
    )

    # Slide 7: wireframe/mock diagram
    s7 = prs.slides[6]
    set_text(s7.shapes[2], "")
    box = s7.shapes[2]
    add_picture_fit(
        s7,
        ROOT / "diagrams" / "wireframe-flow.png",
        box.left,
        box.top,
        box.width,
        box.height,
    )

    # Slide 8: architecture diagram
    s8 = prs.slides[7]
    set_text(s8.shapes[2], "")
    box = s8.shapes[2]
    add_picture_fit(
        s8,
        ROOT / "diagrams" / "architecture.png",
        box.left,
        box.top,
        box.width,
        box.height,
    )

    # Slide 9
    s9 = prs.slides[8]
    set_text(
        s9.shapes[2],
        "Core stack and Google services:\n"
        "- FastAPI + Uvicorn: API orchestration and UI hosting with lightweight deployment.\n"
        "- Google ADK + Vertex AI Gemini 2.5 Flash: structured agent workflows and fast LLM responses.\n"
        "- Vertex runtime bootstrap: auto project/location credential resolution for stable startup.\n"
        "- MCP (FastMCP): standardized tool contracts for YouTube data, analytics, web search, SQLite, Sheets sync.\n"
        "- gRPC services: clean boundaries for YouTube ingestion and Sheets operations.\n"
        "- SQLite (WAL, indexes, retry-safe writes): low-latency local reliability in hackathon constraints.\n"
        "- Google Sheets integration: human-readable operational layer for shared execution tracking.\n"
        "- Tavily search: external context enrichment for strategy quality.\n"
        "- psycopg and phased AlloyDB AI plan: path to scale read-heavy analytics and semantic retrieval.",
        font_size=12,
    )

    # Slide 10: prototype snapshots
    s10 = prs.slides[9]
    set_text(s10.shapes[1], "Prototype snapshots from running system (local instance)", font_size=12)
    clear_text(s10.shapes[3])

    snap_box = s10.shapes[3]
    gap = int(Pt(8))
    cell_w = int((snap_box.width - gap) / 2)
    cell_h = int((snap_box.height - gap) / 2)

    snapshot_items = [
        ("00-login-screen.png", "Login and auth entry"),
        ("01-tasks-dashboard.png", "Tasks workspace with traces"),
        ("02-analytics-tab.png", "Analytics KPIs and insights"),
        ("03-ideas-tab.png", "Ideas and source references"),
    ]

    for idx, (filename, caption) in enumerate(snapshot_items):
        row = idx // 2
        col = idx % 2

        x = int(snap_box.left + col * (cell_w + gap))
        y = int(snap_box.top + row * (cell_h + gap))

        add_picture_fit(
            s10,
            ROOT / "snapshots" / filename,
            x,
            y,
            cell_w,
            int(cell_h - Pt(16)),
        )
        add_textbox(s10, x, int(y + cell_h - Pt(14)), cell_w, int(Pt(14)), caption, font_size=9)

    # Slide 11: close-out summary
    s11 = prs.slides[10]
    add_textbox(
        s11,
        int(Pt(48)),
        int(Pt(120)),
        int(Pt(600)),
        int(Pt(260)),
        "Thank you\n"
        "Q&A\n"
        "Next steps:\n"
        "- Add AlloyDB AI replication for historical runs and semantic analytics.\n"
        "- Add automated evaluation harness for agent output quality.\n"
        "- Harden for deployment with observability and CI tests.",
        font_size=24,
    )

    prs.save(str(PPT_PATH))
    print(f"Updated: {PPT_PATH}")


if __name__ == "__main__":
    main()
